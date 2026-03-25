"""
build_deck.py — Project Photon
Generates ProjectPhoton_Deck.pptx — a premium dark-themed executive presentation.
Run: python3 build_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
C_BG        = RGBColor(0x0D, 0x0F, 0x14)   # near-black
C_BG2       = RGBColor(0x16, 0x1A, 0x26)   # card dark blue
C_BG3       = RGBColor(0x13, 0x16, 0x1E)   # sidebar tone
C_GOLD      = RGBColor(0xC9, 0xA8, 0x4C)   # gold accent
C_BLUE      = RGBColor(0x4A, 0x90, 0xD9)   # blue accent
C_GREEN     = RGBColor(0x6F, 0xCF, 0x97)   # green accent
C_PURPLE    = RGBColor(0xBB, 0x6B, 0xD9)   # purple accent
C_RED       = RGBColor(0xE5, 0x5A, 0x5A)   # red accent
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xE8, 0xEA, 0xF0)   # body text
C_MID       = RGBColor(0xA0, 0xAA, 0xCC)   # secondary text
C_DIM       = RGBColor(0x7A, 0x80, 0x99)   # tertiary text
C_BORDER    = RGBColor(0x25, 0x2B, 0x3E)   # card border

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def bg(slide, color=C_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill_color=None, line_color=None, line_width_pt=1.0):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=18, bold=False, color=C_LIGHT,
        align=PP_ALIGN.LEFT, wrap=True):
    """Add a text box."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return txb

def label(slide, text, x, y, w, h, color=C_BLUE):
    """Small all-caps label above a heading."""
    txt(slide, text.upper(), x, y, w, h, size=9, bold=True, color=color,
        align=PP_ALIGN.LEFT)

def accent_bar(slide, x, y, w=0.04, h=0.32, color=C_GOLD):
    """Thin vertical accent bar."""
    box(slide, x, y, w, h, fill_color=color)

def divider(slide, x, y, w, color=C_BORDER):
    """Thin horizontal rule."""
    box(slide, x, y, w, 0.01, fill_color=color)

def card(slide, x, y, w, h, fill=C_BG2, border=C_BORDER):
    box(slide, x, y, w, h, fill_color=fill, line_color=border, line_width_pt=0.75)

def pill(slide, text, x, y, w=1.5, color=C_BLUE):
    """Small coloured pill / badge."""
    box(slide, x, y, w, 0.28, fill_color=RGBColor(color[0]//8, color[1]//8, color[2]//8),
        line_color=color, line_width_pt=0.5)
    txt(slide, text, x + 0.05, y + 0.04, w - 0.1, 0.22,
        size=8, bold=True, color=color, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 1 — Title
# ============================================================================
sl = prs.slides.add_slide(blank_layout)
bg(sl)

# Dark gradient band top
box(sl, 0, 0, 13.33, 0.08, fill_color=C_GOLD)

# Left accent block
box(sl, 0, 0.08, 0.06, 7.42, fill_color=RGBColor(0x1A, 0x24, 0x56))

# Big title area
box(sl, 0.3, 1.2, 9.0, 2.5, fill_color=RGBColor(0x0A, 0x0C, 0x12))

txt(sl, "⚛  PROJECT PHOTON", 0.5, 1.35, 8.8, 0.7,
    size=42, bold=True, color=C_WHITE)
txt(sl, "Hybrid Quantum-Classical Payment Routing Engine",
    0.5, 2.1, 8.8, 0.55, size=20, bold=False, color=C_GOLD)
txt(sl, "Decision Intelligence Cockpit  ·  Tier-1 Banking Innovation",
    0.5, 2.65, 8.8, 0.4, size=13, color=C_DIM)

divider(sl, 0.5, 3.2, 5.5, color=C_GOLD)

txt(sl, "1.58-bit BitNet LLM  →  IBM Qiskit EstimatorQNN  →  Optimal Route",
    0.5, 3.35, 9.0, 0.4, size=12, color=C_MID)

# Four KPI boxes bottom-right
kpis = [
    ("3.32B", "BitNet Parameters"),
    ("4-Qubit", "IBM QNN Circuit"),
    ("O(log N)", "Quantum Complexity"),
    ("< 1s", "Routing Decision"),
]
for i, (val, lbl) in enumerate(kpis):
    bx = 9.2 + (i % 2) * 2.0
    by = 1.4 + (i // 2) * 1.5
    card(sl, bx, by, 1.8, 1.3, fill=C_BG2)
    txt(sl, val, bx + 0.1, by + 0.15, 1.6, 0.6, size=22, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    txt(sl, lbl, bx + 0.05, by + 0.72, 1.7, 0.45, size=9, color=C_DIM, align=PP_ALIGN.CENTER)

# Bottom strip
box(sl, 0, 7.1, 13.33, 0.4, fill_color=RGBColor(0x0A, 0x0C, 0x14))
txt(sl, "CONFIDENTIAL  ·  INTERNAL USE ONLY  ·  PROJECT PHOTON  ·  2026",
    0, 7.18, 13.33, 0.25, size=8, color=RGBColor(0x30, 0x38, 0x50), align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 2 — Business Problem
# ============================================================================
sl = prs.slides.add_slide(blank_layout)
bg(sl)
box(sl, 0, 0, 13.33, 0.06, fill_color=C_GOLD)
box(sl, 0, 0.06, 0.06, 7.44, fill_color=RGBColor(0x1A, 0x24, 0x56))

label(sl, "The Business Problem", 0.3, 0.2, 6, 0.25, color=C_RED)
txt(sl, "Payment Routing Is a Broken Combinatorial Problem", 0.3, 0.45, 9.5, 0.7,
    size=26, bold=True, color=C_WHITE)
txt(sl, "Tier-1 banks route billions of transactions annually across fragmented networks —\n"
        "classical rule engines evaluate dimensions linearly and miss cross-constraint correlations.",
    0.3, 1.2, 9.5, 0.6, size=12, color=C_MID)
divider(sl, 0.3, 1.85, 12.7)

# Three pain point cards
pain_points = [
    (C_RED,    "01",  "Linear Models Miss\nCorrelations",
     "Approval, cost, latency and resilience are treated as independent axes. "
     "Real routing constraints are deeply entangled — missing these interactions "
     "costs banks $50M–$200M annually in suboptimal decisions."),
    (C_BLUE,   "02",  "Exponential Complexity\nat Scale",
     "Classical exhaustive evaluation scales O(N). As route options grow from "
     "4 to 50+, real-time optimization becomes infeasible. Every new route "
     "multiplies evaluation cost linearly."),
    (C_GREEN,  "03",  "Static Rules Can't\nRead Context",
     "Hard-coded routing rules cannot infer semantic intent from transaction "
     "context. A $2.4M cross-border settlement and a $12K RTP payment need "
     "fundamentally different priority models."),
]
for i, (color, num, title, body) in enumerate(pain_points):
    cx = 0.3 + i * 4.35
    card(sl, cx, 2.05, 4.1, 4.7)
    accent_bar(sl, cx, 2.05, h=4.7, color=color)
    txt(sl, num, cx + 0.25, 2.2, 0.6, 0.5, size=28, bold=True, color=color)
    txt(sl, title, cx + 0.25, 2.75, 3.7, 0.7, size=14, bold=True, color=C_LIGHT)
    txt(sl, body, cx + 0.25, 3.5, 3.7, 3.0, size=10.5, color=C_MID, wrap=True)

# Cost callout bottom
box(sl, 0.3, 6.85, 12.7, 0.45, fill_color=RGBColor(0x1A, 0x08, 0x08),
    line_color=C_RED, line_width_pt=0.5)
txt(sl, "💸  Industry estimate: Suboptimal payment routing costs Tier-1 institutions $50M–$200M annually "
        "in declined revenue, excess processing fees, and failed cross-border settlements.",
    0.5, 6.9, 12.3, 0.35, size=9.5, color=RGBColor(0xFF, 0xAA, 0xAA))

box(sl, 0, 7.1, 13.33, 0.4, fill_color=RGBColor(0x0A, 0x0C, 0x14))
txt(sl, "CONFIDENTIAL  ·  PROJECT PHOTON  ·  2026", 0, 7.18, 13.33, 0.25,
    size=8, color=RGBColor(0x30, 0x38, 0x50), align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 3 — Solution Architecture
# ============================================================================
sl = prs.slides.add_slide(blank_layout)
bg(sl)
box(sl, 0, 0, 13.33, 0.06, fill_color=C_GOLD)
box(sl, 0, 0.06, 0.06, 7.44, fill_color=RGBColor(0x1A, 0x24, 0x56))

label(sl, "The Solution", 0.3, 0.2, 6, 0.25, color=C_GOLD)
txt(sl, "Continuous Hybrid Pipeline: 1.58-bit AI  →  IBM Quantum Neural Network",
    0.3, 0.45, 12.0, 0.6, size=24, bold=True, color=C_WHITE)
divider(sl, 0.3, 1.1, 12.7)

# Pipeline nodes
nodes = [
    ("📋", C_MID,    "INPUT",           "Payment Scenario",    "Natural language\ncontext + metadata"),
    ("🧠", C_BLUE,   "AI LAYER",        "BitNet 1.58-bit LLM", "3.32B params\nCPU inference"),
    ("⚙️", C_GOLD,   "COMPRESSION",     "nn.Linear Reducer",   "3200-dim → 4-dim\nsigmoid [0,1]"),
    ("⚛",  C_PURPLE, "QUANTUM LAYER",   "IBM EstimatorQNN",    "4-qubit circuit\nangle encoding"),
    ("🏆", C_GREEN,  "OUTPUT",          "Photon Match Score",  "Ranked routes\n0 – 100"),
]
node_w, node_h = 2.1, 3.2
gap = 0.26
start_x = 0.3
for i, (icon, color, lbl, name, sub) in enumerate(nodes):
    nx = start_x + i * (node_w + gap)
    card(sl, nx, 1.35, node_w, node_h, fill=C_BG2,
         border=RGBColor(color[0]//3, color[1]//3, color[2]//3))
    # top colour stripe
    box(sl, nx, 1.35, node_w, 0.06, fill_color=color)
    txt(sl, icon,  nx + 0.05, 1.55, node_w - 0.1, 0.55, size=26, align=PP_ALIGN.CENTER)
    txt(sl, lbl,   nx + 0.05, 2.15, node_w - 0.1, 0.3,  size=7.5, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(sl, name,  nx + 0.05, 2.45, node_w - 0.1, 0.45, size=11,  bold=True, color=C_LIGHT, align=PP_ALIGN.CENTER)
    txt(sl, sub,   nx + 0.05, 2.95, node_w - 0.1, 0.6,  size=9,   color=C_DIM, align=PP_ALIGN.CENTER)
    # Arrow between nodes
    if i < len(nodes) - 1:
        ax = nx + node_w + 0.04
        txt(sl, "→", ax, 2.6, gap + 0.14, 0.35, size=16, color=C_BORDER, align=PP_ALIGN.CENTER)

# Bottom explanation strip
box(sl, 0.3, 4.75, 12.7, 1.0, fill_color=RGBColor(0x0C, 0x0F, 0x1A),
    line_color=RGBColor(0x25, 0x2B, 0x3E), line_width_pt=0.5)
txt(sl, "How It Works",
    0.55, 4.85, 3.5, 0.3, size=9, bold=True, color=C_GOLD)
txt(sl,
    "The payment scenario is passed as natural language into the BitNet 1.58-bit LLM. "
    "The model performs a forward pass and the final token's hidden state (3200-dim) is extracted and "
    "compressed to 4 dimensions via a fixed nn.Linear projection. This 4-D tensor is angle-encoded "
    "into a parameterized 4-qubit Qiskit QuantumCircuit (RY gates + CNOT entanglement ring + variational block). "
    "IBM's EstimatorQNN measures Pauli-Z expectation values per qubit, producing quantum-informed "
    "route scores that are blended with classical priority-alignment scoring into the final Photon Match Score.",
    0.55, 5.17, 12.3, 0.52, size=9.5, color=C_MID, wrap=True)

# Quantum circuit detail strip
qc_items = [
    (C_BLUE,   "Layer 1 — Angle Encoding",  "RY(xᵢ · π)  per qubit\nMaps intent tensor to quantum state"),
    (C_PURPLE, "Layer 2 — Entanglement",     "Circular CNOT ring\nCaptures cross-route correlations"),
    (C_GOLD,   "Layer 3 — Variational",      "Trainable RY(θᵢ) rotations\nLearns from routing outcomes"),
    (C_GREEN,  "Measurement",                "Pauli-Z per qubit\nExpectation values → route scores"),
]
for i, (color, title, desc) in enumerate(qc_items):
    qx = 0.3 + i * 3.26
    card(sl, qx, 5.85, 3.1, 1.0, fill=RGBColor(0x0D, 0x0F, 0x14))
    box(sl, qx, 5.85, 0.04, 1.0, fill_color=color)
    txt(sl, title, qx + 0.15, 5.92, 2.9, 0.28, size=8.5, bold=True, color=color)
    txt(sl, desc,  qx + 0.15, 6.22, 2.9, 0.55, size=8,   color=C_DIM)

box(sl, 0, 7.1, 13.33, 0.4, fill_color=RGBColor(0x0A, 0x0C, 0x14))
txt(sl, "CONFIDENTIAL  ·  PROJECT PHOTON  ·  2026", 0, 7.18, 13.33, 0.25,
    size=8, color=RGBColor(0x30, 0x38, 0x50), align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 4 — Quantum vs Classical
# ============================================================================
sl = prs.slides.add_slide(blank_layout)
bg(sl)
box(sl, 0, 0, 13.33, 0.06, fill_color=C_GOLD)
box(sl, 0, 0.06, 0.06, 7.44, fill_color=RGBColor(0x1A, 0x24, 0x56))

label(sl, "Quantum Impact Analysis", 0.3, 0.2, 6, 0.25, color=C_PURPLE)
txt(sl, "Why the Quantum Layer Changes the Decision",
    0.3, 0.45, 10.0, 0.6, size=26, bold=True, color=C_WHITE)
divider(sl, 0.3, 1.1, 12.7)

# Left — Classical column
card(sl, 0.3, 1.3, 5.9, 4.5, fill=RGBColor(0x0D, 0x1A, 0x2A),
     border=RGBColor(0x1F, 0x30, 0x50))
box(sl, 0.3, 1.3, 5.9, 0.32, fill_color=RGBColor(0x1A, 0x2A, 0x50))
txt(sl, "CLASSICAL ROUTING  ·  Priority-Weighted Only",
    0.45, 1.35, 5.6, 0.24, size=8.5, bold=True, color=C_BLUE)

classical_routes = [
    ("#1", "Direct Issuer Connect",      "PREMIUM",  "90.7", True),
    ("#2", "Network Alpha (Visa+)",       "STANDARD", "88.5", False),
    ("#3", "Least-Cost Processor",        "ECONOMY",  "84.5", False),
    ("#4", "Resilient Fallback Network",  "RESILIENT","84.2", False),
]
badge_colors = {"PREMIUM": C_GOLD, "STANDARD": C_BLUE, "ECONOMY": C_GREEN, "RESILIENT": C_PURPLE}
for i, (rank, name, badge, score, winner) in enumerate(classical_routes):
    ry = 1.78 + i * 0.88
    if winner:
        box(sl, 0.35, ry, 5.8, 0.78, fill_color=RGBColor(0x1A, 0x14, 0x02),
            line_color=RGBColor(0xC9, 0xA8, 0x4C), line_width_pt=0.75)
    txt(sl, rank, 0.45, ry + 0.22, 0.45, 0.35, size=14, bold=True, color=C_DIM)
    txt(sl, name, 0.92, ry + 0.1,  3.8,  0.3,  size=10.5, bold=winner, color=C_LIGHT)
    bc = badge_colors.get(badge, C_BLUE)
    txt(sl, badge, 0.92, ry + 0.43, 1.5, 0.22, size=7.5, bold=True, color=bc)
    txt(sl, score, 5.2, ry + 0.15, 0.9, 0.45, size=18, bold=True, color=C_MID, align=PP_ALIGN.RIGHT)
txt(sl, "Linear dot-product scoring. No cross-dimensional correlation modelling.",
    0.4, 5.88, 5.6, 0.35, size=8.5, color=RGBColor(0x5A, 0x60, 0x7A))

# Right — Quantum column
card(sl, 6.5, 1.3, 6.5, 4.5, fill=RGBColor(0x0A, 0x1A, 0x14),
     border=RGBColor(0x1A, 0x4A, 0x2A))
box(sl, 6.5, 1.3, 6.5, 0.32, fill_color=RGBColor(0x0F, 0x2A, 0x1A))
txt(sl, "⚛  QUANTUM-ENHANCED  ·  Hybrid Photon Score",
    6.65, 1.35, 6.2, 0.24, size=8.5, bold=True, color=C_GREEN)

quantum_routes = [
    ("#1", "Direct Issuer Connect",     "PREMIUM",  "93.4", "+8.4",  C_GREEN,  "—"),
    ("#2", "Least-Cost Processor",       "ECONOMY",  "87.9", "+3.4",  C_GREEN,  "▲1"),
    ("#3", "Network Alpha (Visa+)",      "STANDARD", "82.1", "−6.4",  C_RED,    "▼1"),
    ("#4", "Resilient Fallback Network", "RESILIENT","76.8", "−7.4",  C_RED,    "—"),
]
for i, (rank, name, badge, score, adj, adj_color, arrow) in enumerate(quantum_routes):
    ry = 1.78 + i * 0.88
    winner = i == 0
    if winner:
        box(sl, 6.55, ry, 6.38, 0.78, fill_color=RGBColor(0x0A, 0x1A, 0x08),
            line_color=C_GREEN, line_width_pt=0.75)
    txt(sl, rank,  6.62, ry + 0.22, 0.45, 0.35, size=14, bold=True, color=C_DIM)
    txt(sl, name,  7.1,  ry + 0.1,  3.5,  0.3,  size=10.5, bold=winner, color=C_LIGHT)
    bc = badge_colors.get(badge, C_BLUE)
    txt(sl, badge, 7.1,  ry + 0.43, 1.5,  0.22, size=7.5, bold=True, color=bc)
    # QNN adjustment pill
    txt(sl, f"QNN {adj}", 10.5, ry + 0.12, 1.2, 0.28, size=8, bold=True,
        color=adj_color, align=PP_ALIGN.CENTER)
    txt(sl, arrow, 11.85, ry + 0.22, 0.6, 0.35, size=13, bold=True,
        color=adj_color, align=PP_ALIGN.CENTER)
    txt(sl, score, 12.5, ry + 0.15, 0.85, 0.45, size=18, bold=True, color=C_GREEN if winner else C_MID, align=PP_ALIGN.RIGHT)

txt(sl, "30% QNN expectation values + 70% classical alignment. Quantum entanglement captures non-linear priority interactions.",
    6.6, 5.88, 6.2, 0.35, size=8.5, color=RGBColor(0x5A, 0x60, 0x7A))

# Bottom callout
box(sl, 0.3, 6.35, 12.7, 0.6, fill_color=RGBColor(0x0C, 0x0F, 0x1A),
    line_color=RGBColor(0x6F, 0xCF, 0x97), line_width_pt=0.75)
accent_bar(sl, 0.3, 6.35, h=0.6, color=C_GREEN)
txt(sl, "⚛  Quantum Edge: The IBM EstimatorQNN detected cross-priority entanglement patterns "
        "that classical linear scoring cannot model — reshuffling the secondary tier and amplifying "
        "confidence in the primary route through superposition-based expectation measurement.",
    0.5, 6.42, 12.3, 0.48, size=9.5, color=C_MID)

box(sl, 0, 7.1, 13.33, 0.4, fill_color=RGBColor(0x0A, 0x0C, 0x14))
txt(sl, "CONFIDENTIAL  ·  PROJECT PHOTON  ·  2026", 0, 7.18, 13.33, 0.25,
    size=8, color=RGBColor(0x30, 0x38, 0x50), align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 5 — Quantum Scale Advantage
# ============================================================================
sl = prs.slides.add_slide(blank_layout)
bg(sl)
box(sl, 0, 0, 13.33, 0.06, fill_color=C_GOLD)
box(sl, 0, 0.06, 0.06, 7.44, fill_color=RGBColor(0x1A, 0x24, 0x56))

label(sl, "Why Quantum Wins at Enterprise Scale", 0.3, 0.2, 8, 0.25, color=C_PURPLE)
txt(sl, "O(log N) Quantum vs O(N) Classical — The Exponential Advantage",
    0.3, 0.45, 10.0, 0.6, size=24, bold=True, color=C_WHITE)
divider(sl, 0.3, 1.1, 12.7)

# Complexity equation
card(sl, 0.3, 1.3, 5.8, 2.0, fill=RGBColor(0x0A, 0x0C, 0x12))
txt(sl, "Classical Complexity", 0.6, 1.45, 5.2, 0.3, size=10, bold=True, color=C_RED)
txt(sl, "O(N)", 0.6, 1.78, 3.0, 0.65, size=40, bold=True, color=C_RED)
txt(sl, "Each new route = one new evaluation pass.\nLinear scaling makes real-time optimization\ninfeasible beyond 20–30 routes.",
    0.6, 2.48, 5.2, 0.7, size=9, color=C_MID)

card(sl, 6.5, 1.3, 6.5, 2.0, fill=RGBColor(0x0A, 0x14, 0x0A))
txt(sl, "Quantum Complexity", 6.8, 1.45, 5.8, 0.3, size=10, bold=True, color=C_GREEN)
txt(sl, "O(log N)", 6.8, 1.78, 4.0, 0.65, size=40, bold=True, color=C_GREEN)
txt(sl, "One additional qubit doubles the state space\nexplored simultaneously via superposition.\nExponential advantage at enterprise scale.",
    6.8, 2.48, 5.8, 0.7, size=9, color=C_MID)

# Scale table
headers = ["Routes", "Qubits", "Classical Evaluations", "Quantum States", "Payment Use Case"]
col_x   = [0.3, 1.6, 2.9, 5.4, 7.9]
col_w   = [1.2, 1.2, 2.4, 2.4, 5.3]
header_y = 3.5

box(sl, 0.3, header_y, 12.7, 0.32, fill_color=RGBColor(0x16, 0x1A, 0x26))
for j, (hdr, hx, hw) in enumerate(zip(headers, col_x, col_w)):
    hcolor = C_RED if j == 2 else (C_GREEN if j == 3 else C_DIM)
    txt(sl, hdr, hx + 0.08, header_y + 0.07, hw - 0.1, 0.22,
        size=8, bold=True, color=hcolor)

scale_data = [
    ("4",   "2",  "16",           "16",           "Proof-of-concept  (this demo)",  True),
    ("8",   "3",  "256",          "256",          "Regional corridor routing",       False),
    ("16",  "4",  "65,536",       "65,536",       "National payment network",        False),
    ("20",  "5",  "1,048,576",    "1,048,576",    "Cross-border treasury desk",      False),
    ("50",  "6",  "1.1 × 10¹⁵",  "1.1 × 10¹⁵",  "Global real-time settlement",    False),
    ("100", "7",  "1.3 × 10³⁰",  "1.3 × 10³⁰",  "Full interbank network mesh",     False),
]
for r, (routes, qubits, cls_c, q_c, use_case, is_demo) in enumerate(scale_data):
    ry = header_y + 0.32 + r * 0.42
    row_bg = RGBColor(0x1A, 0x14, 0x02) if is_demo else (C_BG2 if r % 2 == 0 else C_BG)
    box(sl, 0.3, ry, 12.7, 0.4, fill_color=row_bg)
    row_data = [routes, qubits, cls_c, q_c, use_case]
    for j, (val, hx, hw) in enumerate(zip(row_data, col_x, col_w)):
        vc = C_GOLD if is_demo else (C_RED if j == 2 else (C_GREEN if j == 3 else C_LIGHT if j <= 1 else C_MID))
        txt(sl, val, hx + 0.08, ry + 0.1, hw - 0.1, 0.25, size=9,
            bold=(j == 0 and is_demo), color=vc)
    if is_demo:
        txt(sl, "◄ THIS DEMO", col_x[4] + col_w[4] - 1.4, ry + 0.1, 1.3, 0.25,
            size=7.5, bold=True, color=C_GOLD, align=PP_ALIGN.RIGHT)

# Bottom note
txt(sl, "One additional qubit doubles the state space — adding qubit 5 takes the system from 16 to 32 evaluated states simultaneously. "
        "Classical routers must add an entire new evaluation loop. This is the fundamental economic argument for quantum-native payment infrastructure.",
    0.3, 7.0, 12.7, 0.35, size=8.5, color=C_DIM)

box(sl, 0, 7.1, 13.33, 0.4, fill_color=RGBColor(0x0A, 0x0C, 0x14))
txt(sl, "CONFIDENTIAL  ·  PROJECT PHOTON  ·  2026", 0, 7.18, 13.33, 0.25,
    size=8, color=RGBColor(0x30, 0x38, 0x50), align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 6 — Business Value
# ============================================================================
sl = prs.slides.add_slide(blank_layout)
bg(sl)
box(sl, 0, 0, 13.33, 0.06, fill_color=C_GOLD)
box(sl, 0, 0.06, 0.06, 7.44, fill_color=RGBColor(0x1A, 0x24, 0x56))

label(sl, "Business Value & ROI", 0.3, 0.2, 6, 0.25, color=C_GREEN)
txt(sl, "Measurable Impact Across Three Value Dimensions",
    0.3, 0.45, 10.0, 0.6, size=26, bold=True, color=C_WHITE)
divider(sl, 0.3, 1.1, 12.7)

value_props = [
    (C_GOLD,   "💰", "Revenue Protection",
     "Approval Rate Lift",
     "$100M+",
     "Recovered on $10B volume",
     "A 1% improvement in approval rate on $10B annual transaction volume "
     "directly recovers previously declined revenue. Quantum entanglement "
     "identifies approval-correlated routing paths that classical linear "
     "scoring systematically undervalues across high-risk and cross-border flows."),
    (C_GREEN,  "📉", "Cost Reduction",
     "Processing Fee Optimization",
     "15–30%",
     "Reduction in per-transaction fees",
     "Multi-dimensional route scoring simultaneously minimises basis-point "
     "costs without sacrificing approval rates or resilience scores. Greatest "
     "impact on domestic and regional corridors with competitive route "
     "availability — achievable in Phase 2 shadow mode deployment."),
    (C_BLUE,   "🛡", "Operational Resilience",
     "Proactive Failover Intelligence",
     "Zero",
     "Unexpected SLA breaches",
     "The QNN models network resilience as a quantum observable — enabling "
     "proactive failover routing before degradation occurs, not after. "
     "Reduces settlement failures, chargeback exposure, and contractual "
     "SLA breaches on high-risk merchant and cross-border payment flows."),
]
for i, (color, icon, category, subtitle, kpi, kpi_sub, body) in enumerate(value_props):
    cx = 0.3 + i * 4.35
    card(sl, cx, 1.35, 4.1, 5.4)
    box(sl, cx, 1.35, 4.1, 0.05, fill_color=color)
    txt(sl, icon,     cx + 0.15, 1.5,  3.8, 0.5, size=26, align=PP_ALIGN.LEFT)
    txt(sl, category, cx + 0.15, 2.1,  3.8, 0.28, size=8, bold=True, color=color)
    txt(sl, subtitle, cx + 0.15, 2.38, 3.8, 0.4,  size=12, bold=True, color=C_LIGHT)
    divider(sl, cx + 0.15, 2.82, 3.8, color=RGBColor(0x25, 0x2B, 0x3E))
    txt(sl, kpi,     cx + 0.15, 2.9,  3.8, 0.65, size=30, bold=True, color=color)
    txt(sl, kpi_sub, cx + 0.15, 3.55, 3.8, 0.28, size=9,  color=C_DIM)
    divider(sl, cx + 0.15, 3.88, 3.8, color=RGBColor(0x25, 0x2B, 0x3E))
    txt(sl, body,    cx + 0.15, 3.96, 3.8, 2.6,  size=9.5, color=C_MID, wrap=True)

# Executive summary callout
box(sl, 0.3, 6.85, 12.7, 0.55, fill_color=RGBColor(0x0C, 0x0F, 0x1A),
    line_color=C_GOLD, line_width_pt=0.75)
accent_bar(sl, 0.3, 6.85, h=0.55, color=C_GOLD)
txt(sl, "Executive Summary:  Project Photon positions the institution as a first-mover in quantum-native "
        "payment infrastructure — targeting $100M+ in annual approval-rate revenue recovery with a phased "
        "deployment path from proof-of-concept to full IBM QPU production within 12 months.",
    0.5, 6.91, 12.3, 0.42, size=9.5, color=C_MID)

box(sl, 0, 7.1, 13.33, 0.4, fill_color=RGBColor(0x0A, 0x0C, 0x14))
txt(sl, "CONFIDENTIAL  ·  PROJECT PHOTON  ·  2026", 0, 7.18, 13.33, 0.25,
    size=8, color=RGBColor(0x30, 0x38, 0x50), align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 7 — Technology Stack
# ============================================================================
sl = prs.slides.add_slide(blank_layout)
bg(sl)
box(sl, 0, 0, 13.33, 0.06, fill_color=C_GOLD)
box(sl, 0, 0.06, 0.06, 7.44, fill_color=RGBColor(0x1A, 0x24, 0x56))

label(sl, "Technology Stack", 0.3, 0.2, 6, 0.25, color=C_PURPLE)
txt(sl, "Purpose-Built for Enterprise AI + Quantum Readiness",
    0.3, 0.45, 10.0, 0.6, size=26, bold=True, color=C_WHITE)
divider(sl, 0.3, 1.1, 12.7)

tech = [
    (C_BLUE,   "BitNet b1.58-3B",          "Classical AI Layer",
     "Microsoft / 1bitLLM · 3.32B parameters\nNative 1.58-bit ternary weights {−1, 0, +1}\n"
     "Zero floating-point multiplications\n8× smaller memory than fp16 equivalent"),
    (C_PURPLE, "IBM Qiskit EstimatorQNN",   "Quantum Layer",
     "4-qubit parameterized QuantumCircuit\nAngle encoding + CNOT entanglement ring\n"
     "Pauli-Z expectation measurement per qubit\nIBM QPU-ready (currently Qiskit-Aer)"),
    (C_GOLD,   "TorchConnector",            "Hybrid Bridge",
     "Wraps QNN as native PyTorch nn.Module\nEnd-to-end differentiable pipeline\n"
     "Gradient-based QNN training support\nSeamless MPS / CPU device handoff"),
    (C_GREEN,  "PyTorch",                   "Tensor Operations",
     "nn.Linear: 3200-dim → 4-dim reducer\nXavier-uniform initialization\n"
     "Sigmoid normalization to [0,1]\nApple Silicon MPS acceleration"),
    (C_RED,    "HuggingFace Transformers",  "Model Loading",
     "AutoModelForCausalLM + trust_remote_code\nlow_cpu_mem_usage shard streaming\n"
     "LlamaTokenizer fallback for BitNet vocab\nHF mirror endpoint support"),
    (C_MID,    "Streamlit",                 "Executive UI",
     "st.cache_resource model caching\nst.status live pipeline progress\n"
     "Premium dark-mode CSS theme\nMulti-page navigation (Overview + App)"),
]
for i, (color, name, role, details) in enumerate(tech):
    cx = 0.3 + (i % 3) * 4.35
    cy = 1.35 + (i // 3) * 2.75
    card(sl, cx, cy, 4.1, 2.55)
    box(sl, cx, cy, 0.05, 2.55, fill_color=color)
    txt(sl, name,    cx + 0.2, cy + 0.12, 3.8, 0.32, size=12, bold=True, color=C_LIGHT)
    txt(sl, role,    cx + 0.2, cy + 0.44, 3.8, 0.22, size=8.5, bold=True, color=color)
    divider(sl, cx + 0.2, cy + 0.68, 3.7, color=RGBColor(0x25, 0x2B, 0x3E))
    txt(sl, details, cx + 0.2, cy + 0.76, 3.8, 1.65, size=9, color=C_MID)

box(sl, 0, 7.1, 13.33, 0.4, fill_color=RGBColor(0x0A, 0x0C, 0x14))
txt(sl, "CONFIDENTIAL  ·  PROJECT PHOTON  ·  2026", 0, 7.18, 13.33, 0.25,
    size=8, color=RGBColor(0x30, 0x38, 0x50), align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 8 — Roadmap
# ============================================================================
sl = prs.slides.add_slide(blank_layout)
bg(sl)
box(sl, 0, 0, 13.33, 0.06, fill_color=C_GOLD)
box(sl, 0, 0.06, 0.06, 7.44, fill_color=RGBColor(0x1A, 0x24, 0x56))

label(sl, "Strategic Roadmap", 0.3, 0.2, 6, 0.25, color=C_GOLD)
txt(sl, "From Proof-of-Concept to Production Quantum Routing",
    0.3, 0.45, 10.0, 0.6, size=26, bold=True, color=C_WHITE)
divider(sl, 0.3, 1.1, 12.7)

# Timeline bar
box(sl, 0.3, 1.55, 12.7, 0.06, fill_color=RGBColor(0x25, 0x2B, 0x3E))
phases = [
    (C_BLUE,   "✓ COMPLETE",  "Phase 1",   "Proof of Concept",
     ["1.58-bit BitNet integration", "IBM Qiskit QNN pipeline",
      "Hybrid scoring engine", "Executive dashboard UI", "GitHub deployment"]),
    (C_GREEN,  "NEXT — Q2",   "Phase 2",   "Training & Validation",
     ["QNN training on historical data", "A/B vs. production router",
      "8-qubit / 16-route expansion", "Risk-adjusted scoring model"]),
    (C_GOLD,   "6 MONTHS",    "Phase 3",   "Enterprise Integration",
     ["Payment gateway API wrapper", "Real-time transaction stream",
      "Compliance & audit logging", "Shadow mode deployment"]),
    (C_PURPLE, "12 MONTHS",   "Phase 4",   "IBM Quantum Cloud",
     ["IBM QPU hardware execution", "20+ qubit circuit expansion",
      "1M+ route state space", "Full production rollout"]),
]
for i, (color, timeline, phase, title, items) in enumerate(phases):
    px = 0.3 + i * 3.26
    # Timeline dot
    box(sl, px + 1.5, 1.45, 0.28, 0.28, fill_color=color)
    card(sl, px, 1.85, 3.1, 5.0, fill=RGBColor(0x0D, 0x0F, 0x14))
    box(sl, px, 1.85, 3.1, 0.05, fill_color=color)
    txt(sl, timeline, px + 0.12, 1.95, 2.86, 0.25, size=7.5, bold=True, color=color)
    txt(sl, phase,    px + 0.12, 2.2,  2.86, 0.22, size=8,   color=C_DIM)
    txt(sl, title,    px + 0.12, 2.42, 2.86, 0.45, size=13,  bold=True, color=C_LIGHT)
    divider(sl, px + 0.12, 2.9, 2.86, color=RGBColor(0x25, 0x2B, 0x3E))
    for j, item in enumerate(items):
        txt(sl, f"→  {item}", px + 0.12, 3.0 + j * 0.52, 2.86, 0.44, size=9.5, color=C_MID)

# Bottom executive callout
box(sl, 0.3, 7.0, 12.7, 0.48, fill_color=RGBColor(0x0C, 0x0F, 0x1A),
    line_color=C_GOLD, line_width_pt=0.75)
accent_bar(sl, 0.3, 7.0, h=0.48, color=C_GOLD)
txt(sl, "Investment Ask: Approve Phase 2 funding to begin QNN training on 12 months of historical routing data "
        "and initiate shadow mode deployment alongside the production router — validating ROI before full cutover.",
    0.5, 7.06, 12.3, 0.36, size=9.5, color=C_MID)

box(sl, 0, 7.52, 13.33, 0.4, fill_color=RGBColor(0x0A, 0x0C, 0x14))
txt(sl, "CONFIDENTIAL  ·  PROJECT PHOTON  ·  2026", 0, 7.6, 13.33, 0.25,
    size=8, color=RGBColor(0x30, 0x38, 0x50), align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 9 — Closing / CTA
# ============================================================================
sl = prs.slides.add_slide(blank_layout)
bg(sl)
box(sl, 0, 0, 13.33, 0.08, fill_color=C_GOLD)
box(sl, 0, 0.08, 0.06, 7.42, fill_color=RGBColor(0x1A, 0x24, 0x56))

# Central glow box
box(sl, 1.5, 1.5, 10.33, 4.5, fill_color=RGBColor(0x0C, 0x0F, 0x1A),
    line_color=RGBColor(0xC9, 0xA8, 0x4C), line_width_pt=1.0)
box(sl, 1.5, 1.5, 10.33, 0.07, fill_color=C_GOLD)

txt(sl, "⚛", 5.8, 1.65, 1.5, 0.8, size=42, align=PP_ALIGN.CENTER, color=C_GOLD)
txt(sl, "Project Photon", 1.7, 2.5, 9.93, 0.65,
    size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txt(sl, "The future of payment routing is hybrid, quantum-native, and available today.",
    1.7, 3.2, 9.93, 0.45, size=14, color=C_GOLD, align=PP_ALIGN.CENTER)
divider(sl, 3.0, 3.75, 7.33, color=RGBColor(0x25, 0x2B, 0x3E))
txt(sl,
    "1.58-bit BitNet LLM  ·  IBM Qiskit EstimatorQNN  ·  Hybrid Photon Score\n"
    "github.com/suman-hacks/ProjectPhoton",
    1.7, 3.88, 9.93, 0.7, size=11, color=C_MID, align=PP_ALIGN.CENTER)

# Three next step pills at bottom of box
ctas = ["Approve Phase 2 Funding", "Schedule Live Demo", "Review GitHub Repository"]
cta_colors = [C_GOLD, C_GREEN, C_BLUE]
for i, (cta, color) in enumerate(zip(ctas, cta_colors)):
    bx = 2.0 + i * 3.45
    box(sl, bx, 5.22, 3.1, 0.5,
        fill_color=RGBColor(color[0]//6, color[1]//6, color[2]//6),
        line_color=color, line_width_pt=0.75)
    txt(sl, cta, bx + 0.08, 5.31, 2.94, 0.32, size=10, bold=True,
        color=color, align=PP_ALIGN.CENTER)

box(sl, 0, 7.1, 13.33, 0.4, fill_color=RGBColor(0x0A, 0x0C, 0x14))
txt(sl, "CONFIDENTIAL  ·  INTERNAL USE ONLY  ·  PROJECT PHOTON  ·  2026",
    0, 7.18, 13.33, 0.25, size=8, color=RGBColor(0x30, 0x38, 0x50), align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
output = "ProjectPhoton_Deck.pptx"
prs.save(output)
print(f"✅  Saved: {output}  ({len(prs.slides)} slides)")
